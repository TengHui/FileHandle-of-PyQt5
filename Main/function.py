import os
import glob
import shutil
from pdf2docx import Converter
from PyPDF2 import PdfFileReader, PdfFileWriter
import win32com.client
from docx import Document
from pptx import Presentation
from docx2pdf import convert
import filecmp


class FileHandle:
    '''
    构造函数
    source_pathlist:输入的文件夹/文件绝对路径列表
    (可选参数)
    goal_path:输入的目标文件夹绝对路径
    split_list:输入分割序列
    password:密码
    '''

    def __init__(self, source_pathlist, goal_path='', split_list=[], password=''):
        self.source_pathlist = source_pathlist
        if not '/' in goal_path and not '\\' in goal_path:
            #默认为第一个文件的同级目录下
            self.goal_path=os.path.split(source_pathlist[0])[0]
        else:
            self.goal_path = goal_path
        self.split_list = split_list
        self.password = password
        self.file_list = []
        self.FileList()

    '''
    self.source_pathlist处理函数
    功能:辨别输入的路径是文件夹还是文件,处理成纯文件路径
    '''

    def FileList(self):
        self.file_list = []
        for sourcefile in self.source_pathlist:
            if os.path.exists(sourcefile):
                if os.path.isfile(sourcefile):
                    self.file_list.append(sourcefile)
                if os.path.isdir(sourcefile):
                    self.GetFileName(sourcefile)

    '''
    文件夹处理函数:获取文件夹下所有文件绝对路径
    '''
    '''
    def GetFileName(self,sourcefile):
        # 文件路径拼接,获取绝对路径
        filelist = [os.path.join(dirpath, filesname) \
                    for dirpath, dirs, files in os.walk(sourcefile) \
                    for filesname in files]
        return filelist
    '''
    def GetFileName(self,rootdir):
        #根目录+子文件夹+文件
        for root,dirs,files in os.walk(rootdir):
            for file in files:
                file_path=os.path.join(root,file)
                self.file_list.append(file_path)
            for dir in dirs:
                dir_path=os.path.join(root,dir)
                self.GetFileName(dir_path)#递归调用

    '''
    PDF模块功能函数
    1.pdf->word:PdfToWord();
    2.pdf合并:PdfMerge();
    3.pdf分割:PdfSplit();
    4.pdf加密:PdfEncrypte();
    '''

    # 1.pdf->word
    def PdfToWord(self):
        for pdf_file in self.file_list:
            # 对文件名称后缀进行判断
            suffix = pdf_file.split('/')[-1].split('.')[-1]
            if suffix == 'pdf':
                docx_file = self.goal_path + '/' + pdf_file.split('/')[-1].split('.')[0] + '.docx'
                cv = Converter(pdf_file)
                cv.convert(docx_file)
                cv.close()

    # 2.pdf合并
    def PdfMerge(self):
        # 写入器
        pdf_writer = PdfFileWriter()
        for pdf_file in self.file_list:
            # 对文件名称后缀进行判断
            suffix = pdf_file.split('/')[-1].split('.')[-1]
            if suffix == 'pdf':
                # 读取器
                pdf_reader = PdfFileReader(open(pdf_file, 'rb'))
                pageCount = pdf_reader.getNumPages()
                for page in range(pageCount):
                    pdf_writer.addPage(pdf_reader.getPage(page))
        # 将写入器内的pdf写入输出文件
        with open(os.path.join(self.goal_path, "合并.pdf"), 'wb') as outputfile:
            pdf_writer.write(outputfile)

    # 3.pdf分割
    def PdfSplit(self):
        for pdf_file in self.file_list:
            # 对文件名称后缀进行判断
            suffix = pdf_file.split('/')[-1].split('.')[-1]
            if suffix == 'pdf':
                pdf_reader = PdfFileReader(open(pdf_file, 'rb'))
                # 获取文件名称和路径
                output_path, pdf = os.path.splitext(pdf_file)
                # 如果不存在目标文件夹,再新建文件夹
                if not os.path.exists(output_path):
                    os.mkdir(output_path)
                startPage = 0  # 分割起始页
                endPage = 0  # 分割终止页
                for page in self.split_list:
                    pdf_writer = PdfFileWriter()
                    # 更新终止页
                    endPage += page
                    for iPage in range(startPage, endPage):
                        pdf_writer.addPage(pdf_reader.getPage(iPage))
                    # 拼接存储路径
                    save_path = output_path + '/' + pdf_file.split('/')[-1].split('.')[
                        0] + f'({startPage + 1}-{endPage}).pdf'
                    with open(save_path, 'wb') as outputfile:
                        pdf_writer.write(outputfile)
                    # 更新起始页
                    startPage += page

    # 4.pdf加密
    def PdfEncrypte(self):
        for pdf_file in self.file_list:
            # 对文件名称后缀进行判断
            suffix = pdf_file.split('/')[-1].split('.')[-1]
            if suffix == 'pdf':
                pdf_reader = PdfFileReader(pdf_file)
                pdf_writer = PdfFileWriter()
                for page in range(pdf_reader.getNumPages()):
                    pdf_writer.addPage(pdf_reader.getPage(page))
                # 给写入器添加密码
                pdf_writer.encrypt(self.password)
                # 默认给文件添加密码,源文件将被直接添加密码，不在产生新文件
                with open(pdf_file, 'wb') as outputfile:
                    pdf_writer.write(outputfile)

    '''
    PPT模块功能函数
    1.PPT->pdf:PptToPdf();
    2.PPT->word:PptToWord();
    3.PPT->image:PptToImage();
    '''

    # 1.PPT->pdf
    def PptToPdf(self):
        # 打开ppt应用程序
        ppt_app = win32com.client.Dispatch('PowerPoint.Application')
        for ppt_file in self.file_list:
            # 对文件名称后缀进行判断
            suffix = ppt_file.split('/')[-1].split('.')[-1]
            if suffix == 'ppt' or suffix == 'pptx':
                ppt = ppt_app.Presentations.Open(ppt_file)
                save_path = self.goal_path + '/' + ppt_file.split('/')[-1].split('.')[0] + '.pdf'
                save_path=save_path.replace('/', '\\')
                print(save_path)
                ppt.SaveAs(save_path, 32)
                ppt.Close()
        # 关闭应用程序
        ppt_app.Quit()

    # 2.PPT->word:只能提取文本框的文字内容
    def PptToWord(self):
        for ppt_file in self.file_list:
            # 对文件名称后缀进行判断
            suffix = ppt_file.split('/')[-1].split('.')[-1]
            if suffix == 'ppt' or suffix == 'pptx':
                word_file = Document()
                pptx = Presentation(ppt_file)
                # 遍历所有文本框
                for slide in pptx.slides:
                    for shape in slide.shapes:
                        if shape.has_text_frame:
                            text_frame = shape.text_frame
                            for paragraph in text_frame.paragraphs:
                                word_file.add_paragraph(paragraph.text)
                save_path = self.goal_path + '/' + ppt_file.split('/')[-1].split('.')[0] + '.docx'
                word_file.save(save_path)

    # 3.PPT->image:固定jpg格式
    def PptToImage(self):
        # 打开ppt应用程序
        ppt_app = win32com.client.Dispatch('PowerPoint.Application')
        for ppt_file in self.file_list:
            # 对文件名称后缀进行判断
            suffix = ppt_file.split('/')[-1].split('.')[-1]
            if suffix == 'ppt' or suffix == 'pptx':
                ppt = ppt_app.Presentations.Open(ppt_file)
                save_path = self.goal_path + '/' + ppt_file.split('/')[-1].split('.')[0] + '.pdf'
                save_path=save_path.replace('/', '\\')
                print(save_path)
                ppt.SaveAs(save_path, 17)
                ppt.Close()
        # 关闭应用程序
        ppt_app.Quit()

    '''
    word模块功能函数
    1.word->pdf:WordToPdf();
    2.word合并:WordMerge();
    3.word分割:WordSplit();
    '''

    # 1.word->pdf
    def WordToPdf(self):
        for word_file in self.file_list:
            # 对文件名称后缀进行判断
            suffix = word_file.split('/')[-1].split('.')[-1]
            if suffix == 'doc' or suffix == 'docx':
                convert(word_file, self.goal_path)

    # 2.word合并
    def WordMerge(self):
        # 启动word应用对象
        word = win32com.client.Dispatch('Word.Application')
        # 应用程序窗口不显示
        word.Visible = False
        # 建立空文档
        output_file = word.Documents.Add()
        # 拼接文档
        for word_file in self.file_list:
            # 对文件名称后缀进行判断
            suffix = word_file.split('/')[-1].split('.')[-1]
            if suffix == 'doc' or suffix == 'docx':
                output_file.Application.Selection.InsertFile(word_file)
        # 获取合并后文档内容
        doc = output_file.Range(output_file.Content.Start, output_file.Content.End)
        # 莫名其妙只能存为doc格式
        save_path = self.goal_path + '/' + '合并.doc'
        output_file.SaveAs(save_path)
        output_file.Close()

    # 3.word分割
    def WordSplit(self):
        pass

    '''
    文件处理模块功能函数
    1.文件整理:FileArrange();
    2.文件去重:FileDuplicate();
    '''

    # 1.文件整理:按文件后缀名称整理，存入对应文件夹
    def FileArrange(self):
        # 如果不存在目标文件夹,则建立
        save_path=self.goal_path+'/'+'分类'
        if not os.path.exists(save_path):
            os.mkdir(save_path)
        for source_dir in self.source_pathlist:
            #for file in glob.glob(f'{source_dir}/**/*', recursive=True):
            for file in self.file_list:
                if os.path.isfile(file):
                    filename = os.path.basename(file)
                    if '.' in filename:
                        suffix = filename.split('.')[-1]
                    else:
                        suffix = 'others'
                    # 建立对应文件后缀名称的文件夹
                    if not os.path.exists(f'{save_path}/{suffix}'):
                        os.mkdir(f'{save_path}/{suffix}')
                    # 复制文件过去:默认情况
                    shutil.copy(file, f'{save_path}/{suffix}')
                    # 移动文件
                    #shutil.move(file, f'{self.goal_path}/{suffix}')

    # 2.文件去重:即使文件名称不同,也会比较文件内容进行去重
    def FileDuplicate(self):
        '''
        # 深层次文件列表
        filelist = []
        # 遍历搜索文件夹，建立文件列表
        for source_dir in self.source_pathlist:
            for i in glob.glob(source_dir + '/**/*', recursive=True):
                if os.path.isfile(i):
                    filelist.append(i)
        '''
        # 循环比较所有文件
        for f1 in self.file_list:
            for f2 in self.file_list:
                if f1 != f2 and os.path.exists(f1) and os.path.exists(f2):
                    if filecmp.cmp(f1, f2):
                        os.remove(f2)
